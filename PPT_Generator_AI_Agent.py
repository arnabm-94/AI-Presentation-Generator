from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0]  # Title slide layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Artificial Intelligence & Emerging Technologies"
subtitle.text = "AI, ML, Generative AI, Vertex AI & Copilot vs ChatGPT"

# Slide 2: Artificial Intelligence Overview
slide_layout = prs.slide_layouts[1]  # Title and content layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Artificial Intelligence (AI)"
content.text = "- AI refers to the simulation of human intelligence in machines.\n" \
               "- Applications: Automation, Data Analysis, Predictive Modeling.\n" \
               "- Impact: Transforms industries like healthcare, finance, and manufacturing."

# Slide 3: Machine Learning Overview
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Machine Learning (ML)"
content.text = "- Subset of AI that enables machines to learn from data.\n" \
               "- Types: Supervised, Unsupervised, Reinforcement Learning.\n" \
               "- Applications: Fraud Detection, Recommendation Systems, Image Recognition."

# Slide 4: Generative AI
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Generative AI"
content.text = "- AI that generates text, images, and videos.\n" \
               "- Uses neural networks like GANs and transformers.\n" \
               "- Applications: Content creation, AI art, Chatbots, Code generation."

# Slide 5: Vertex AI
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Vertex AI"
content.text = "- Google's AI platform for building and deploying ML models.\n" \
               "- Features: AutoML, custom model training, scalable infrastructure.\n" \
               "- Benefits: Cost-efficient, seamless integration with Google Cloud."

# Slide 6: Why Copilot is Better Than ChatGPT?
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Why Copilot is Better Than ChatGPT?"
content.text = "- Copilot integrates with Microsoft tools, enhancing productivity.\n" \
               "- ChatGPT focuses on conversational AI but lacks deep integration.\n" \
               "- Copilot provides contextual assistance in real-time applications."

# Slide 7: Conclusion
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Conclusion"
content.text = "- AI and ML are transforming industries with advanced automation.\n" \
               "- Generative AI is revolutionizing content creation.\n" \
               "- Vertex AI offers powerful ML model deployment.\n" \
               "- Copilot provides seamless AI assistance for productivity."

# Save the presentation
pptx_filename = "C:/Users/ARNMUKHE/OneDrive - Capgemini/ARNAB MUKHERJEE/AI & AUTOMATION PROJECTS/Automatic PPT Generator/AI_Technologies_Presentation.pptx"
prs.save(pptx_filename)

pptx_filename


